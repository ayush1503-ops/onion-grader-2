"""Build the SIH26031 pitch deck (6 slides) as .pptx using python-pptx."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BROWN = RGBColor(0x7A, 0x3B, 0x12)
GOLD = RGBColor(0xD9, 0x8E, 0x3F)
DARK = RGBColor(0x2A, 0x1C, 0x10)
LIGHT = RGBColor(0xFF, 0xF6, 0xEA)
GREY = RGBColor(0x66, 0x5A, 0x4A)
GREEN = RGBColor(0x1A, 0x7A, 0x1A)
RED = RGBColor(0xC0, 0x39, 0x2B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def bg(slide, color=LIGHT):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def title_bar(slide, text):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(1.0))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(34); p.font.bold = True; p.font.color.rgb = BROWN

def body(slide, lines, top=1.5, left=0.6, width=12.1, height=5.6, size=20):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame; tf.word_wrap = True
    for i, (txt, bold, col, sz) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(sz or size)
        p.font.bold = bold
        p.font.color.rgb = col or DARK
        p.space_after = Pt(10)
    return box

def band(slide, y=6.7, text="", color=GOLD):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(y), prs.slide_width, Inches(0.8))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    tf = sh.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = LIGHT

# ---- Slide 1: Title ----
s = prs.slides.add_slide(BLANK); bg(s, BROWN)
b = s.shapes.add_textbox(Inches(1), Inches(2.1), Inches(11.3), Inches(2.6))
tf = b.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "OnionQuality: AI-based Onion Grading App"
p.font.size = Pt(44); p.font.bold = True; p.font.color.rgb = LIGHT; p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph(); p2.text = "Smart India Hackathon 2026  ·  SIH26031"
p2.font.size = Pt(24); p2.font.color.rgb = GOLD; p2.alignment = PP_ALIGN.CENTER
p3 = tf.add_paragraph(); p3.text = "Ministry of Consumer Affairs, Food & Public Distribution — DoCA"
p3.font.size = Pt(18); p3.font.color.rgb = RGBColor(0xF0, 0xE2, 0xCE); p3.alignment = PP_ALIGN.CENTER
b2 = s.shapes.add_textbox(Inches(1), Inches(5.4), Inches(11.3), Inches(1.2))
tf = b2.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "Team <Your Team Name>"
p.font.size = Pt(20); p.font.color.rgb = LIGHT; p.alignment = PP_ALIGN.CENTER

# ---- Slide 2: Problem ----
s = prs.slides.add_slide(BLANK); bg(s)
title_bar(s, "The Problem")
body(s, [
 ("Onion quality grading at procurement centres is SUBJECTIVE —", True, BROWN, 24),
 ("it depends on the person looking at the lot, and it varies from centre to centre.", False, DARK, 22),
 ("", False, DARK, 12),
 ("This causes:", True, GREY, 20),
 ("• disputes between farmers and procurement agencies", False, DARK, 22),
 ("• inconsistent prices for the same quality of onions", False, DARK, 22),
 ("• no instant, shareable record of quality at the time of procurement", False, DARK, 22),
 ("", False, DARK, 12),
 ("A fair, fast, transparent, AI-assisted grading system is needed.", True, RED, 24),
])

# ---- Slide 3: Solution ----
s = prs.slides.add_slide(BLANK); bg(s)
title_bar(s, "Our Solution — one photo in, a full report out")
body(s, [
 ("A mobile web app. Take ONE photo of the batch → get instant quality grading.", True, BROWN, 24),
 ("", False, DARK, 10),
 ("1.  DETECTS every onion in the photo (computer vision)", False, DARK, 20),
 ("2.  CLASSIFIES each: Damaged / Rotten / Sprouted / Undersized / Good", False, DARK, 20),
 ("3.  MEASURES diameter (mm) using a coin as scale reference", False, DARK, 20),
 ("4.  COMPUTES Grade A % and URS % (relaxed-spec) automatically", False, DARK, 20),
 ("5.  GENERATES a digital quality report instantly (JSON + card image)", False, DARK, 20),
])

# ---- Slide 4: How it works ----
s = prs.slides.add_slide(BLANK); bg(s)
title_bar(s, "How it works — the pipeline")
body(s, [
 ("[ Photo ] → [ Segment each onion ] → [ Extract features ] → [ Classify ] → [ Grade % ] → [ Report ]", True, GREY, 18),
 ("", False, DARK, 8),
 ("Segmentation:  Otsu thresholding + morphology → find each onion's contour", False, DARK, 20),
 ("Features (visible only):", True, BROWN, 20),
 ("   • Colour (HSV) — brown/dark = damage or rot, green = sprout", False, DARK, 20),
 ("   • Shape — circularity, size (diameter in mm via coin reference)", False, DARK, 20),
 ("   • Texture — surface roughness via gray-level statistics", False, DARK, 20),
 ("Rules → GOOD (Grade A 45–65 mm / URS 35–70 mm) or REJECT", False, DARK, 20),
])

# ---- Slide 5: Honest limitations ----
s = prs.slides.add_slide(BLANK); bg(s)
title_bar(s, "What a camera can and cannot see (we are honest)")
body(s, [
 ("CAN detect (visible surface):", True, GREEN, 22),
 ("   discoloration, bruises, rot/mold on the skin, sprouts, undersized bulbs, shape", False, DARK, 20),
 ("", False, DARK, 8),
 ("CANNOT detect:", True, RED, 22),
 ("   internal rot, hidden damage, internal moisture — anything not visible in the photo", False, DARK, 20),
 ("", False, DARK, 8),
 ("This boundary is stated in the app and the report — transparency is the point.", True, BROWN, 22),
])

# ---- Slide 6: Roadmap ----
s = prs.slides.add_slide(BLANK); bg(s)
title_bar(s, "Roadmap")
body(s, [
 ("Phase 1 (built):  rule-based grader + mobile web app + digital report", False, DARK, 20),
 ("Phase 2:  labelled dataset → ML classifier → better accuracy per defect type", False, DARK, 20),
 ("Phase 3:  CNN on-device (TensorFlow Lite) for offline Android/iOS app", False, DARK, 20),
 ("Phase 4:  regional-language UI (Hindi + others), procurement-centre integration", False, DARK, 20),
 ("", False, DARK, 10),
 ("Live demo — next.", True, BROWN, 26),
])
band(s, y=6.7, text="Live demo follows", color=GOLD)

prs.save("/home/user/onion-opencv-course/SIH26031_pitch_deck.pptx")
print("saved pitch deck, slides:", len(prs.slides._sldIdLst))
